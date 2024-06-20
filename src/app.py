import chainlit as cl
from openai import OpenAI, AzureOpenAI
from pptx import Presentation
from docx import Document
import PyPDF2
import os
import json
import time
import random
from types import SimpleNamespace

# Get environment variables
OPENAI_KEY =  os.getenv("OPENAI_KEY")
AZURE_OPENAI_KEY =  os.getenv("AZURE_OPENAI_KEY")
AZURE_OPENAI_ENDPOINT =  os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_VERSION =  os.getenv("AZURE_OPENAI_VERSION")
GPT_MODEL = os.getenv("GPT_MODEL")
TEMPERATURE = float(os.getenv("TEMPERATURE"))
with open("system_prompt.txt", 'r', encoding='utf-8') as file:
    SYSTEM_PROMPT = file.read()

TEMPLATE_FILE = os.getenv("TEMPLATE_FILE")
TITLE_TEMPLATE_SLIDE_INDEX = int(os.getenv("TITLE_TEMPLATE_SLIDE_INDEX"))
CONTENT_TEMPLATE_SLIDE_INDEX = int(os.getenv("CONTENT_TEMPLATE_SLIDE_INDEX"))

TEMP_FILES_FOLDER = ".files"
WAITING_MESSAGE = "Please wait..."

if len(OPENAI_KEY) > 0:
    ai_client = OpenAI(api_key = OPENAI_KEY)
elif len(AZURE_OPENAI_KEY) > 0:
    ai_client = AzureOpenAI(api_key = AZURE_OPENAI_KEY, api_version=AZURE_OPENAI_VERSION, azure_endpoint=AZURE_OPENAI_ENDPOINT)
else:
    print("[ERROR] Need to set up API key for OpenAI or Azure OpenAI")
    exit(1)

tools = [
    {
        "type": "function",
        "function": {
            "name": "generate_presentation",
            "description": "Generate powerpoint presentation slides",
            "parameters": {
                "type": "object",
                "properties": {
                    "topic": {
                        "type": "string",
                        "description": "Topic of the presentation",
                    },
                    "slide_data": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "Title of the slide"
                                },
                                "content": {
                                    "type": "array",
                                    "items": {
                                        "type": "string",
                                        "description": "Content for one bullet point"
                                    },
                                    "description": "An array of main contents of the slide"
                                }
                            }
                        },
                        "description": "An array of slide contents",
                    }
                },
                "required": ["topic", "slide_data"],
            }
        }
    }
]

@cl.on_message
async def main(message: cl.Message):
    # Create folder for temporary files for current user if not exist
    user_temp_file_folder = f'{TEMP_FILES_FOLDER}/{cl.user_session.get("id")}'
    if not os.path.exists(user_temp_file_folder):
        os.makedirs(user_temp_file_folder)
    
    save_message_to_history(message)

    #Create an empty response to show loading icon
    response_msg = cl.Message(content="")
    await response_msg.send()
    
    #Call OpenAI
    conversation_history = get_conversation_history()
    result = await cl.make_async(get_gpt_response)(ai_client, GPT_MODEL, TEMPERATURE, SYSTEM_PROMPT, conversation_history, tools)
    if result.content:
        response = result.content
    elif result.tool_calls:
        function_name = result.tool_calls[0].function.name
        arguments = json.loads(result.tool_calls[0].function.arguments)
        if function_name == "generate_presentation":
            try:
                topic = arguments["topic"]
                slide_data = arguments["slide_data"]
                generated_file_path = await cl.make_async(create_powerpoint_file)(topic, slide_data, f'{TEMP_FILES_FOLDER}/{cl.user_session.get("id")}')
                attached_file = cl.File(name=topic, path=generated_file_path, display="inline")
                response_msg.elements = [attached_file]
                response = f'[SUCCESS] Your Powerpoint presentation has been generated successfully'
            except Exception as e:             
                response = f'[ERROR] Problem generating org chart:\n {e}'
        else:
            response = f"[ERROR] Invalid function"
    else:
        response = f"[ERROR] Invalid response from OpenAI"
    
    response_msg.content = response
    await response_msg.update()
    save_message_to_history(response_msg)

#============================================#

def save_message_to_history(message):
    if cl.user_session.get("chat_history"):
        chat_history = cl.user_session.get("chat_history")
    else:
        chat_history = []
    chat_history.append(message)
    cl.user_session.set("chat_history", chat_history)

def get_conversation_history():
    result = []
    if cl.user_session.get("chat_history"):
        for message in cl.user_session.get("chat_history"):
            if message.author == "User":
                processed_message = {"role": "user", "content": message.content}
                if message.elements:
                    try:
                        attached_files = [file for file in message.elements]
                        first_file = attached_files[0]
                        file_extension = get_file_extension(first_file.name).lower()
                        if file_extension in [".txt", ".docx", ".pdf"]:
                            file_content = read_all_text_from_file(first_file.path, file_extension)
                            prompt_content = message.content + ":\nHere is the provided information in the attached document:\n" + file_content
                            processed_message = {"role": "user", "content": prompt_content}
                    except Exception as e:             
                        pass
                result.append(processed_message)
            else:
                result.append({"role": "assistant", "content": message.content})
    return result

def generate_random_file_name():
    return f'{int(time.time_ns())}_{random.randint(0,10000)}'

def get_file_extension(file_path):
    _, file_extension = os.path.splitext(file_path)
    return file_extension

def read_all_text_from_file(file_path, file_extension = None):
    result = ""
    if not file_extension:
        file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == '.txt' or file_extension == '.text':
        with open(file_path, 'r', encoding='utf-8') as file:
            result = file.read()
    elif file_extension == '.docx':
        document = Document(file_path)
        result = '\n'.join([para.text for para in document.paragraphs])
    elif file_extension == '.pdf':
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                result += page.extract_text() or ''     
    return result


def create_powerpoint_file(topic, slide_data, output_folder):    
    # Load the template
    prs = Presentation(TEMPLATE_FILE)

    # Add title slide
    slide_layout = prs.slide_layouts[TITLE_TEMPLATE_SLIDE_INDEX]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "Generated by Slides Assistant"
    
    # Add content slides
    for slide_content in slide_data:
        new_slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_TEMPLATE_SLIDE_INDEX])
        if "title" in slide_content:
            new_slide.shapes.title.text = slide_content["title"]
        if "content" in slide_content:
            new_slide.placeholders[1].text = "\n".join(slide_content["content"])
    
    # Save the presentation
    output_file_path = f'{output_folder}/{generate_random_file_name()}.pptx'
    prs.save(output_file_path)
    return output_file_path

def get_gpt_response(ai_client, gpt_model, temperature, system_prompt, conversation_history, tools):
    prompt_structure = [{"role": "system", "content": system_prompt}]
    for msg in conversation_history:
        prompt_structure.append(msg) 
    try:
        response = ai_client.chat.completions.create(
            model = gpt_model,
            messages = prompt_structure,
            temperature = temperature,
            tools = tools,
            tool_choice = "auto"
        )
        return response.choices[0].message
    except Exception as e:
        return SimpleNamespace(content=f"[ERROR] Problem calling OpenAI API:\n {e}")